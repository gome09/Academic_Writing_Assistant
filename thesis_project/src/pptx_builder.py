# -*- coding: utf-8 -*-
"""
PPT 生成器 —— 按 PPT_SPEC 生成答辩演示草案 (.pptx)。

已落实的规范：
  - 16:9 画布 (13.333 x 7.5 in)
  - 字号层级：封面 40 / 章节 32 / 页标题 30 / 正文 24
  - 白底 + 主色标题条；左对齐；每页要点 <=6
  - 结构：封面 -> 目录 -> 4 个分节 + 内容页 -> 致谢
  - 页脚含页码与论文题目
  - 中文字体：微软雅黑（含 eastAsia 设置）
"""
from __future__ import annotations
import io
import math
import sys
import os
import unicodedata

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from config.format_spec import PPT_SPEC as P

_C = P["theme"]
PRIMARY = RGBColor(*_C["primary_rgb"])
ACCENT = RGBColor(*_C["accent_rgb"])
TEXT = RGBColor(*_C["text_rgb"])
MUTED = RGBColor(*_C["muted_rgb"])
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LOGGER = None

CN_FONT = P["font"]["family_cn"]
EN_FONT = P["font"]["family_en"]
SZ = P["sizes"]
W_IN = P["slide"]["width_inch"]
H_IN = P["slide"]["height_inch"]


class BuildResult(str):
    """build() 返回值：既是输出路径字符串，又携带本次构建的告警列表（T0-4）。

    作为 str 子类保持完全向后兼容——os.path.getsize / endswith / 字符串拼接
    等调用方无需感知 .warnings 属性即可正常工作。
    """
    warnings: list

    def __new__(cls, path: str, warnings: list):
        obj = str.__new__(cls, path)
        obj.warnings = warnings
        return obj


def reload_spec():
    """Refresh cached presentation constants after applying a YAML template."""
    global _C, PRIMARY, ACCENT, TEXT, MUTED, CN_FONT, EN_FONT, SZ, W_IN, H_IN
    _C = P["theme"]
    PRIMARY = RGBColor(*_C["primary_rgb"])
    ACCENT = RGBColor(*_C["accent_rgb"])
    TEXT = RGBColor(*_C["text_rgb"])
    MUTED = RGBColor(*_C["muted_rgb"])
    CN_FONT = P["font"]["family_cn"]
    EN_FONT = P["font"]["family_en"]
    SZ = P["sizes"]
    W_IN = P["slide"]["width_inch"]
    H_IN = P["slide"]["height_inch"]


def _set_font(run, size_pt, color=None, bold=False, cn=None):
    color = color or TEXT
    cn = cn or CN_FONT
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = EN_FONT
    # eastAsia 中文字体
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", cn)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 全空白版式


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _warn(message, warnings: list):
    """记录告警到局部列表并打印（T0-4：不再依赖模块级 LAST_WARNINGS）。"""
    warnings.append(message)
    if _LOGGER is not None:
        _LOGGER.warning(message)
    print(f"  [警告] {message}")


def _textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def _rect(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    _fill(shp, color)
    shp.shadow.inherit = False
    return shp


# ---------------------------------------------------------------------------
#  页脚（页码 + 题目）
# ---------------------------------------------------------------------------
def _footer(slide, idx, total, title):
    tf = _textbox(slide, 0.5, H_IN - 0.45, W_IN - 1.0, 0.35)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title[:30]
    _set_font(r, 10, MUTED)
    tf2 = _textbox(slide, W_IN - 1.5, H_IN - 0.45, 1.0, 0.35)
    pp = tf2.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    rr = pp.add_run(); rr.text = f"{idx} / {total}"
    _set_font(rr, 10, MUTED)


# ---------------------------------------------------------------------------
#  各类幻灯片
# ---------------------------------------------------------------------------
def _slide_cover(prs, s):
    slide = _blank(prs)
    _rect(slide, 0, 0, W_IN, H_IN, WHITE)
    _rect(slide, 0, H_IN * 0.38, W_IN, 0.06, PRIMARY)   # 装饰细条
    tf = _textbox(slide, 1.0, 2.2, W_IN - 2.0, 2.0, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = s["title"]
    _set_font(r, SZ["cover_title_pt"], PRIMARY, bold=True)
    tf2 = _textbox(slide, 1.0, 4.6, W_IN - 2.0, 1.5)
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = s.get("subtitle", "")
    _set_font(r, SZ["body_pt"], TEXT)
    p2 = tf2.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = "本科毕业论文答辩"
    _set_font(r, SZ["caption_pt"], MUTED)
    return slide


def _slide_outline(prs, s):
    slide = _blank(prs)
    _title_bar(slide, s["title"])
    tf = _textbox(slide, 1.2, 1.8, W_IN - 2.4, H_IN - 2.5)
    for i, item in enumerate(s["items"], 1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.space_after = Pt(14)
        r = p.add_run(); r.text = f"0{i}   {item}"
        _set_font(r, SZ["section_title_pt"] - 4, TEXT)
    return slide


def _slide_section(prs, s):
    slide = _blank(prs)
    _rect(slide, 0, 0, W_IN, H_IN, PRIMARY)
    tf = _textbox(slide, 1.0, H_IN / 2 - 0.8, W_IN - 2.0, 1.6, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = s["title"]
    _set_font(r, SZ["section_title_pt"], WHITE, bold=True)
    return slide


def _slide_content(prs, s, warnings=None):
    if warnings is None:
        warnings = []
    slide = _blank(prs)
    _title_bar(slide, s["title"])
    media = s.get("media") or []
    has_media = bool(media)
    max_content = W_IN * P["layout"].get("content_max_ratio", 0.66)
    text_width = 6.0 if has_media else min(W_IN - 2.0, max_content)
    tf = _textbox(slide, 1.0, 1.8, text_width, H_IN - 2.6)
    bullets = s["bullets"][:P["layout"]["max_bullets_per_slide"]]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        align = P["layout"].get("text_align", "left")
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        p.space_after = Pt(12)
        r = p.add_run(); r.text = "▪  " + b
        _set_font(r, SZ["body_pt"], TEXT)
    if has_media:
        _render_media(slide, media[0], 7.2, 1.9, 5.4, 4.8, warnings)
    return slide


def _render_media(slide, media, left, top, width, height, warnings=None):
    """Render the first media block using deterministic, bounded layouts."""
    if warnings is None:
        warnings = []
    if media.get("kind") == "image" and media.get("data"):
        try:
            slide.shapes.add_picture(io.BytesIO(media["data"]),
                                     Inches(left), Inches(top),
                                     width=Inches(width), height=Inches(height))
        except Exception as exc:  # noqa: BLE001
            _warn(f"PPT 图片插入失败：{exc}", warnings)
            tf = _textbox(slide, left, top + height / 2 - 0.2, width, 0.5)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = "<图片插入失败>"
            _set_font(r, SZ["caption_pt"], MUTED)
        return
    if media.get("kind") != "table":
        return
    rows = media.get("rows") or []
    if not rows:
        return
    max_rows = P["layout"].get("table_max_rows", 8)
    max_cols = P["layout"].get("table_max_cols", 6)
    rows = rows[:max_rows]
    n_cols = min(max(len(r) for r in rows), max_cols)
    shape = slide.shapes.add_table(len(rows), n_cols, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    table = shape.table
    for ri, row in enumerate(rows):
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            cell.text = str(row[ci]) if ci < len(row) else ""
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _set_font(run, 12, WHITE if ri == 0 else TEXT,
                              bold=(ri == 0))
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY


def _slide_thanks(prs, s):
    slide = _blank(prs)
    _rect(slide, 0, 0, W_IN, H_IN, WHITE)
    tf = _textbox(slide, 1.0, H_IN / 2 - 1.0, W_IN - 2.0, 2.0, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = s["title"]
    _set_font(r, SZ["cover_title_pt"], PRIMARY, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = s.get("subtitle", "")
    _set_font(r, SZ["body_pt"], TEXT)
    return slide


def _title_bar(slide, title):
    """页标题：左侧主色竖条 + 标题文字。"""
    _rect(slide, 0.5, 0.55, 0.12, 0.7, ACCENT)
    tf = _textbox(slide, 0.8, 0.5, W_IN - 1.6, 0.9, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set_font(r, SZ["slide_title_pt"], PRIMARY, bold=True)
    # 标题下细分隔线
    _rect(slide, 0.5, 1.45, W_IN - 1.0, 0.02, RGBColor(0xDD, 0xDD, 0xDD))


# ---------------------------------------------------------------------------
#  主构建
# ---------------------------------------------------------------------------
_DISPATCH = {
    "cover": _slide_cover,
    "outline": _slide_outline,
    "section": _slide_section,
    "content": _slide_content,
    "thanks": _slide_thanks,
}


def build(deck, out_path):
    global _LOGGER
    try:
        import logging
        _LOGGER = logging.getLogger("thesis_project")
    except Exception:  # pragma: no cover - logging is part of stdlib
        _LOGGER = None
    reload_spec()
    warnings: list = []  # T0-4：局部告警列表，随返回值传出
    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)

    slides = deck["slides"]
    total = len(slides)
    for idx, s in enumerate(slides, 1):
        fn = _DISPATCH.get(s["type"], _slide_content)
        # content 页可能含媒体，需要传入 warnings；其余页不产生告警
        if s["type"] == "content":
            slide = fn(prs, s, warnings)
        else:
            slide = fn(prs, s)
        # 封面/分节/致谢不加页脚
        if s["type"] in ("outline", "content"):
            _footer(slide, idx, total, deck["title"])
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]

    p_min = P["principle"]["total_slides_min"]
    p_max = P["principle"]["total_slides_max"]
    if total < p_min:
        _warn(f"PPT 总页数 {total} 低于规范下限 {p_min}，请检查。", warnings)
    elif total > p_max:
        _warn(f"PPT 总页数 {total} 高于规范上限 {p_max}，请检查。", warnings)
    _check_structure(slides, warnings)
    _check_content_limits(slides, warnings)
    _check_shape_bounds(prs, warnings)
    prs.save(out_path)
    return BuildResult(out_path, warnings)


def _check_structure(slides, warnings=None):
    """按 PPT_SPEC["structure"] 校验各段页数。

    - cover / outline / thanks：按 slide["type"] 直接计数校验。
      其中 outline 只校验上限：现有合法 deck（见 tests/test_pptx_page_count.py
      的构造方式）允许省略目录页，缺失目录不作警告。
    - background / method / result / conclusion 四段：content 页按
      slide["bucket"] 归属计数，与规范 min/max 比较（section 分节页不计入）。
      整份 deck 都不带 bucket 信息（旧格式/手工构造）时跳过该项校验。
    """
    if warnings is None:
        warnings = []
    checkable = {"cover", "outline", "thanks"}
    skip_min = {"outline"}
    counts = {}
    for s in slides:
        counts[s["type"]] = counts.get(s["type"], 0) + 1
    for seg in P["structure"]:
        key = seg["key"]
        if key not in checkable:
            continue
        n = counts.get(key, 0)
        if n < seg["min"] and key not in skip_min:
            _warn(f"PPT {seg['title']}页数 {n} 低于规范下限 {seg['min']}，请检查。",
                  warnings)
        elif n > seg["max"]:
            _warn(f"PPT {seg['title']}页数 {n} 高于规范上限 {seg['max']}，请检查。",
                  warnings)

    seg_by_key = {seg["key"]: seg for seg in P["structure"]}
    bucket_counts = {}
    for s in slides:
        if s.get("type") == "content" and s.get("bucket"):
            bucket_counts[s["bucket"]] = bucket_counts.get(s["bucket"], 0) + 1
    if not bucket_counts:
        return  # 旧格式/手工 deck 不带 bucket 信息，跳过分段页数校验
    for key in ("background", "method", "result", "conclusion"):
        seg = seg_by_key.get(key)
        if seg is None:
            continue
        n = bucket_counts.get(key, 0)
        if n < seg["min"]:
            _warn(f"PPT {seg['title']}内容页数 {n} "
                  f"低于规范下限 {seg['min']}，请检查。", warnings)
        elif n > seg["max"]:
            _warn(f"PPT {seg['title']}内容页数 {n} "
                  f"高于规范上限 {seg['max']}，请检查。", warnings)


def _display_width(text):
    return sum(2 if unicodedata.east_asian_width(ch) in "WFA" else 1
               for ch in str(text))


def _check_content_limits(slides, warnings=None):
    """Warn about content that cannot satisfy configured line/table limits."""
    if warnings is None:
        warnings = []
    max_lines = P["layout"].get("max_lines_per_bullet", 2)
    for idx, slide in enumerate(slides, 1):
        if slide.get("type") != "content":
            continue
        chars_per_line = P["layout"].get("chars_per_line_media", 26) \
            if slide.get("media") else P["layout"].get("chars_per_line_text", 44)
        for bullet in slide.get("bullets", []):
            if math.ceil(_display_width(bullet) / chars_per_line) > max_lines:
                _warn(f"PPT 第 {idx} 页要点可能超过 {max_lines} 行："
                      f"{str(bullet)[:30]}", warnings)
        for media in slide.get("media", []):
            if media.get("kind") == "table":
                cols = max((len(r) for r in media.get("rows") or []), default=0)
                max_cols = P["layout"].get("table_max_cols", 6)
                if cols > max_cols:
                    _warn(f"PPT 第 {idx} 页表格有 {cols} 列，"
                          f"PPT 仅展示前 {max_cols} 列，请在 Word 中查看完整表格。",
                          warnings)


def _check_shape_bounds(prs, warnings=None):
    if warnings is None:
        warnings = []
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if (shape.left < 0 or shape.top < 0 or
                    shape.left + shape.width > prs.slide_width or
                    shape.top + shape.height > prs.slide_height):
                _warn(f"PPT 第 {idx} 页存在越界元素："
                      f"{getattr(shape, 'name', 'shape')}", warnings)
