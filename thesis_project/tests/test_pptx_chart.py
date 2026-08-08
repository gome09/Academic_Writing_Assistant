# -*- coding: utf-8 -*-
"""T2-6：PPT 原生图表 + AI 配图占位测试。

不依赖文生图，保持本地可用；表格可图表化时渲染 python-pptx 原生柱状图。
"""
from __future__ import annotations

from pptx import Presentation

from config.format_spec import PPT_SPEC
from src import pptx_builder


# ---------------------------------------------------------------------------
#  _try_float / _table_is_chartable 纯函数
# ---------------------------------------------------------------------------
def test_try_float_parses_numbers():
    assert pptx_builder._try_float("123") == 123.0
    assert pptx_builder._try_float("1,234.5") == 1234.5
    assert pptx_builder._try_float("87%") == 87.0
    assert pptx_builder._try_float(42) == 42.0


def test_try_float_rejects_non_numeric():
    assert pptx_builder._try_float("abc") is None
    assert pptx_builder._try_float("") is None
    assert pptx_builder._try_float(None) is None


def test_table_is_chartable_numeric_column():
    rows = [["季度", "销售额"], ["Q1", "100"], ["Q2", "200"], ["Q3", "300"]]
    assert pptx_builder._table_is_chartable(rows) is True


def test_table_not_chartable_text_only():
    rows = [["方法", "优点"], ["A", "简单"], ["B", "准确"]]
    assert pptx_builder._table_is_chartable(rows) is False


def test_table_not_chartable_too_few_rows():
    rows = [["季度", "销售额"], ["Q1", "100"]]
    assert pptx_builder._table_is_chartable(rows) is False


def test_table_not_chartable_single_column():
    rows = [["标题"], ["a"], ["b"], ["c"]]
    assert pptx_builder._table_is_chartable(rows) is False


# ---------------------------------------------------------------------------
#  _render_chart 生成原生图表
# ---------------------------------------------------------------------------
def _blank_slide():
    prs = Presentation()
    prs.slide_width = pptx_builder.Inches(pptx_builder.W_IN)
    prs.slide_height = pptx_builder.Inches(pptx_builder.H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def test_render_chart_produces_chart_shape():
    rows = [["季度", "销售额"], ["Q1", "100"], ["Q2", "200"], ["Q3", "300"]]
    prs, slide = _blank_slide()
    ok = pptx_builder._render_chart(slide, rows, 1.0, 1.0, 5.0, 4.0)
    assert ok is True
    charts = [s for s in slide.shapes if s.has_chart]
    assert len(charts) == 1
    # 系列名取自表头
    assert charts[0].chart.series[0].name == "销售额"


def test_render_chart_returns_false_for_non_chartable():
    rows = [["方法", "优点"], ["A", "简单"], ["B", "准确"]]
    prs, slide = _blank_slide()
    ok = pptx_builder._render_chart(slide, rows, 1.0, 1.0, 5.0, 4.0)
    assert ok is False
    assert not any(s.has_chart for s in slide.shapes)


# ---------------------------------------------------------------------------
#  集成：build() 在 chart_from_table 开启时把可图表化表格渲染为图表
# ---------------------------------------------------------------------------
def _deck_with_table_media(rows):
    """构造单页 content deck，媒体为给定表格。"""
    return {"title": "T", "slides": [
        {"type": "cover", "title": "题目", "subtitle": ""},
        {"type": "content", "title": "数据展示",
         "bullets": ["要点一"],
         "media": [{"kind": "table", "rows": rows}]},
        {"type": "thanks", "title": "谢谢", "subtitle": ""},
    ]}


def _count_shapes(out_path):
    """统计生成文件中的图表与表格形状数。"""
    prs = Presentation(str(out_path))
    charts = tables = 0
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_chart:
                charts += 1
            if shp.has_table:
                tables += 1
    return charts, tables


def test_build_renders_chart_when_enabled(tmp_path, monkeypatch):
    """chart_from_table=True 时，可图表化表格渲染为原生图表而非表格。"""
    monkeypatch.setitem(PPT_SPEC["layout"], "chart_from_table", True)
    rows = [["季度", "销售额"], ["Q1", "100"], ["Q2", "200"], ["Q3", "300"]]
    out = tmp_path / "chart.pptx"
    pptx_builder.build(_deck_with_table_media(rows), str(out))
    charts, tables = _count_shapes(out)
    assert charts == 1
    assert tables == 0


def test_build_keeps_table_when_disabled(tmp_path):
    """默认 chart_from_table=False，表格仍渲染为表格。"""
    assert PPT_SPEC["layout"]["chart_from_table"] is False
    rows = [["季度", "销售额"], ["Q1", "100"], ["Q2", "200"], ["Q3", "300"]]
    out = tmp_path / "table.pptx"
    pptx_builder.build(_deck_with_table_media(rows), str(out))
    charts, tables = _count_shapes(out)
    assert charts == 0
    assert tables == 1


def test_build_keeps_nonchartable_table_as_table(tmp_path, monkeypatch):
    """chart_from_table=True 但表格不可图表化时，仍渲染为表格。"""
    monkeypatch.setitem(PPT_SPEC["layout"], "chart_from_table", True)
    rows = [["方法", "优点"], ["A", "简单"], ["B", "准确"]]
    out = tmp_path / "text_table.pptx"
    pptx_builder.build(_deck_with_table_media(rows), str(out))
    charts, tables = _count_shapes(out)
    assert charts == 0
    assert tables == 1


# ---------------------------------------------------------------------------
#  AI 配图占位（image_placeholder）
# ---------------------------------------------------------------------------
def _deck_content_no_media():
    return {"title": "T", "slides": [
        {"type": "cover", "title": "题目", "subtitle": ""},
        {"type": "content", "title": "研究方法",
         "bullets": ["要点一", "要点二"]},
        {"type": "thanks", "title": "谢谢", "subtitle": ""},
    ]}


def _slide_texts(out_path):
    """收集所有幻灯片的文本内容。"""
    prs = Presentation(str(out_path))
    texts = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame:
                texts.append(shp.text_frame.text)
    return texts


def test_image_placeholder_added_when_enabled(tmp_path, monkeypatch):
    """image_placeholder=True 时，无媒体内容页出现配图占位框与图题建议。"""
    monkeypatch.setitem(PPT_SPEC["layout"], "image_placeholder", True)
    out = tmp_path / "placeholder.pptx"
    pptx_builder.build(_deck_content_no_media(), str(out))
    texts = _slide_texts(out)
    joined = "\n".join(texts)
    assert "配图占位" in joined
    assert "研究方法" in joined  # 图题建议含幻灯片标题


def test_image_placeholder_absent_when_disabled(tmp_path):
    """默认 image_placeholder=False，无媒体内容页不出现占位框。"""
    assert PPT_SPEC["layout"]["image_placeholder"] is False
    out = tmp_path / "no_placeholder.pptx"
    pptx_builder.build(_deck_content_no_media(), str(out))
    texts = _slide_texts(out)
    assert "配图占位" not in "\n".join(texts)
