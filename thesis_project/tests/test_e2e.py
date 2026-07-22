# -*- coding: utf-8 -*-
import os

from src import llm_enhancer, docx_builder, pptx_builder
from src.readers import read_dir
from src.organizer import organize

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


def test_enhance_skipped_without_key(monkeypatch):
    monkeypatch.delenv("LLM_API_KEY", raising=False)
    docs = read_dir(os.path.join(ROOT, "sample_input"))
    thesis, deck = organize(docs)
    t2, d2 = llm_enhancer.enhance(thesis, deck, docs)
    assert t2 is thesis and d2 is deck  # 原样返回


def test_enhance_survives_step_failures(monkeypatch):
    monkeypatch.setenv("LLM_API_KEY", "sk-test")

    def boom(s, u):
        raise RuntimeError("api down")
    monkeypatch.setattr(llm_enhancer, "_chat_json", boom)
    docs = read_dir(os.path.join(ROOT, "sample_input"))
    thesis, deck = organize(docs)
    t2, d2 = llm_enhancer.enhance(thesis, deck, docs)
    # 全部步骤失败也要拿回可用的草案数据
    assert t2["title"] == "基于深度学习的校园垃圾图像分类系统设计与实现"
    assert any(s["type"] == "content" for s in d2["slides"])


def test_full_pipeline_sample(tmp_path):
    docs = read_dir(os.path.join(ROOT, "sample_input"))
    thesis, deck = organize(docs)
    assert thesis["title"] == "基于深度学习的校园垃圾图像分类系统设计与实现"
    assert thesis["auto_skeleton"] is False
    wp = docx_builder.build(thesis, str(tmp_path / "论文草案.docx"))
    pp = pptx_builder.build(deck, str(tmp_path / "答辩PPT草案.pptx"))
    assert os.path.getsize(wp) > 0
    assert os.path.getsize(pp) > 0

    # 内容级断言：产物不仅存在，还要包含关键结构
    import docx as _docx
    from pptx import Presentation as _P

    d = _docx.Document(wp)
    all_text = "\n".join(p.text for p in d.paragraphs)
    assert "参考文献" in all_text
    assert "[1]" in all_text                      # 参考文献编号
    assert "关键词" in all_text
    prs = _P(pp)
    cover_texts = [sh.text_frame.text for sh in prs.slides[0].shapes
                   if sh.has_text_frame]
    assert any(t.strip() for t in cover_texts)    # 封面非空
    assert len(prs.slides) >= 5
