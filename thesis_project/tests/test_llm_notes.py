# -*- coding: utf-8 -*-
from pptx import Presentation

from src import llm_enhancer, pptx_builder


def test_add_speaker_notes_attaches_marked_notes(monkeypatch):
    monkeypatch.setattr(llm_enhancer, "_chat_json",
                        lambda s, u: {"绪论": "各位老师好，本章介绍研究背景。"})
    deck = {"title": "t", "slides": [
        {"type": "cover", "title": "t", "subtitle": "s"},
        {"type": "content", "title": "绪论", "bullets": ["背景"]},
    ]}
    llm_enhancer.add_speaker_notes(deck)
    assert deck["slides"][1]["notes"].startswith(llm_enhancer.AI_MARK)
    assert "研究背景" in deck["slides"][1]["notes"]
    assert "notes" not in deck["slides"][0]          # 非 content 页不写


def test_add_speaker_notes_failure_is_silent(monkeypatch):
    def boom(s, u):
        raise RuntimeError("超时")
    monkeypatch.setattr(llm_enhancer, "_chat_json", boom)
    deck = {"title": "t", "slides": [
        {"type": "content", "title": "绪论", "bullets": ["a"]}]}
    llm_enhancer.add_speaker_notes(deck)             # 不抛异常
    assert "notes" not in deck["slides"][0]


def test_pptx_builder_writes_notes(tmp_path):
    deck = {"title": "题目", "slides": [
        {"type": "cover", "title": "题目", "subtitle": "答辩人：张三"},
        {"type": "content", "title": "绪论", "bullets": ["背景"],
         "notes": "这里是口播备注。"},
    ]}
    out = str(tmp_path / "o.pptx")
    pptx_builder.build(deck, out)
    prs = Presentation(out)
    assert prs.slides[1].notes_slide.notes_text_frame.text == "这里是口播备注。"


# ---------------------------------------------------------------------------
#  T2-5: 演讲备注含时长估算
# ---------------------------------------------------------------------------
def test_notes_include_duration_estimation(monkeypatch):
    """T2-5：备注含【约X秒】时长标注。"""
    monkeypatch.setattr(llm_enhancer, "_chat_json",
                        lambda s, u: {"绪论": "各位老师好。", "方法": "本章介绍设计。"})
    deck = {"title": "t", "slides": [
        {"type": "cover", "title": "t", "subtitle": "s"},
        {"type": "content", "title": "绪论", "bullets": ["背景"]},
        {"type": "content", "title": "方法", "bullets": ["设计"]},
    ]}
    llm_enhancer.add_speaker_notes(deck, talk_minutes=10)
    for s in deck["slides"][1:]:
        assert "【约" in s["notes"]
        assert "秒】" in s["notes"]


def test_estimate_duration_basic():
    """10分钟(600秒)/3页 -> 每页约200秒。"""
    dur = llm_enhancer._estimate_duration("title", 3, 10)
    assert "200" in dur


def test_estimate_duration_minimum():
    """时长不低于30秒。"""
    dur = llm_enhancer._estimate_duration("title", 100, 10)
    assert "30" in dur


def test_estimate_duration_zero():
    """0分钟或0页 -> 空字符串。"""
    assert llm_enhancer._estimate_duration("t", 0, 10) == ""
    assert llm_enhancer._estimate_duration("t", 10, 0) == ""
